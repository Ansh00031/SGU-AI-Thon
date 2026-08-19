"""Smart PowerPoint generator and template updater.

Searches for existing PPTX files in common directories (Downloads, Desktop, etc.),
injects all slide contents, or generates a complete 16:9 presentation file.
"""

import os
import glob
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("[*] Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

from generate_ppt import create_deck

SLIDE_DATA = {
    1: {
        "title": "Autonomous OS Debugging Agent",
        "subtitle": "AI-Powered System Diagnostics, Safe Remediation & Zero-Risk Rollback",
        "team_name": "KernelHealers",
        "problem_statement": "Autonomous Diagnosis & Remediation of Operating System Errors using AI Agents",
        "members": "Ansh & Team",
        "college": "Your College Name",
    },
    2: {
        "title": "PROBLEM STATEMENT",
        "points": [
            "• Cryptic Error Codes: Codes like 0x80070005 (Access Denied) or 0x80240020 leave users with zero actionable guidance.",
            "• Dangerous Forum Copy-Pasting: Users spend 45+ mins copy-pasting unverified commands from Reddit/StackOverflow that risk corrupting system registries.",
            "• Massive Support Overhead: Over 30% of enterprise IT helpdesk time is wasted on repetitive permission and background service failures.",
        ]
    },
    3: {
        "title": "PROPOSED SOLUTION",
        "points": [
            "• Live Machine Ingestion: Directly queries Windows Event Viewer logs (System, App, WindowsUpdate) & OS telemetry in real-time.",
            "• Closed-Loop AI Reasoning: Formulates hypotheses and executes safe, read-only system probes (icacls, Get-Service) to confirm root cause.",
            "• Human-in-the-Loop Safety Gate: Generates syntax-highlighted PowerShell scripts requiring explicit user approval ([Y/n]) before execution.",
            "• Pre-Fix Snapshot & Instant Rollback: Automatically captures system restore points before remediation, enabling 1-click reversal ('agent.py rollback').",
            "• Privacy-First & Local LLM Support: Runs seamlessly with cloud GPT-4o or 100% offline via local Ollama models (Llama 3.1).",
        ]
    },
    4: {
        "title": "FLOW OF SOLUTION",
        "points": [
            "1. Ingestion: Privilege checks + Windows Event Viewer critical logs (50 events) + OS telemetry extraction.",
            "2. AI Diagnosis: Formulates diagnostic theory & generates safe read-only PowerShell commands.",
            "3. System Probing: Executes read-only probes (icacls, service queries) with command regex blacklists.",
            "4. Human Approval: Syntax-highlighted script preview with strict [Y/n] confirmation gate.",
            "5. Sandboxed Fix & Verification: Pre-fix snapshot creation, isolated execution, post-fix health check, and instant rollback support.",
        ]
    },
    5: {
        "title": "TECH STACK & IMPLEMENTATION",
        "points": [
            "• Core CLI & UI: Python 3.10+, Typer (Type-safe CLI framework), Rich (Formatted banners, tables, and Monokai syntax highlighting).",
            "• AI Engine: OpenAI API (GPT-4o) / Local Ollama (Llama 3.1 / Mistral) with strict JSON output schemas.",
            "• OS Diagnostics: PowerShell 5.1/7+, Windows Event Log API (Get-WinEvent), Win32 / WMI system APIs.",
            "• Safety Guardrails: Regex command blacklists, subprocess sandbox with auto-cleanup, and session snapshot manager (.backups/).",
        ]
    },
    6: {
        "title": "UNIQUE SELLING PROPOSITION (USP)",
        "points": [
            "• Ground Truth Telemetry (Zero Hallucinations): Runs live read-only verification probes on actual files/services before suggesting fixes.",
            "• Strict Human-in-the-Loop Transparency: Full preview of the exact PowerShell script with mandatory affirmative confirmation.",
            "• Zero-Risk Snapshot Rollback: Every fix automatically creates an inverse rollback script and snapshot point (revert in under 1s).",
            "• Air-Gapped / Enterprise Ready: Runs 100% offline on local LLMs for high-security, defense, or banking computers.",
        ]
    },
    7: {
        "title": "FEASIBILITY & COMPETITOR ANALYSIS",
        "points": [
            "• Windows Troubleshooter: Basic & rigid, black-box fixes that usually fail on complex permission errors.",
            "• ChatGPT / Web Search: Requires manual copy-pasting, lacks machine context, and risks running dangerous commands.",
            "• Enterprise RMM (Tanium/Datadog): Extremely expensive, requires complex setup and manual script authoring.",
            "• Our AI Debug Agent: Automated deep diagnosis, 1-click execution with approval, 1-click rollback, free & open-source.",
        ]
    },
    8: {
        "title": "RESEARCH, REFERENCES & FUTURE SCOPE",
        "points": [
            "• GitHub Repository: github.com/[your-username]/os-debug-agent (MIT Open Source License).",
            "• Live Validation: Successfully diagnosed and remediated Windows Update permission denial (0x80070005) & service failure modes.",
            "• Future Scope: Multi-OS Linux (systemd/journalctl) & macOS support, fleet management mode for DevOps, and interactive CLI chat mode.",
        ]
    }
}


def find_template_and_update():
    # Search common directories for PPTX files
    search_paths = [
        Path.cwd(),
        Path.home() / "Downloads",
        Path.home() / "Desktop",
        Path.home() / "Documents",
        Path(r"C:\Users\ansh6\.gemini\antigravity\scratch"),
    ]

    target_template = None
    for p in search_paths:
        if p.exists():
            for f in p.glob("*.pptx"):
                if "bharat" in f.name.lower() or "hackathon" in f.name.lower() or "build" in f.name.lower():
                    target_template = f
                    break
        if target_template:
            break

    if target_template and target_template.exists():
        print(f"[*] Found existing PPTX template: {target_template}")
        try:
            prs = Presentation(str(target_template))
            for slide_num, slide in enumerate(prs.slides, 1):
                if slide_num in SLIDE_DATA:
                    data = SLIDE_DATA[slide_num]
                    # Find or add text
                    # Add content text box
                    left = Inches(1.0)
                    top = Inches(1.8)
                    width = Inches(11.33)
                    height = Inches(4.8)
                    tx_box = slide.shapes.add_textbox(left, top, width, height)
                    tf = tx_box.text_frame
                    tf.word_wrap = True

                    if slide_num == 1:
                        p = tf.paragraphs[0]
                        p.text = f"Project: {data['title']}\n{data['subtitle']}\n\nTeam: {data['team_name']}\nMembers: {data['members']}\nCollege: {data['college']}"
                        p.font.size = Pt(18)
                        p.font.bold = True
                    else:
                        for idx, pt in enumerate(data.get("points", [])):
                            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                            p.text = pt
                            p.font.size = Pt(14)
                            p.font.name = "Arial"
                            p.space_after = Pt(12)

            out_path = Path.cwd() / "Build_With_Bharat_Final_Presentation.pptx"
            prs.save(str(out_path))
            print(f"[+] Updated template and saved to: {out_path}")
            return
        except Exception as ex:
            print(f"[!] Error modifying existing template: {ex}. Generating standalone presentation...")

    # Fallback to generating the high quality presentation directly
    out_path = Path.cwd() / "Build_With_Bharat_Final_Presentation.pptx"
    create_deck(str(out_path))


if __name__ == "__main__":
    find_template_and_update()
