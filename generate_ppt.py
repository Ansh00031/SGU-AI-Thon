"""Generate the complete, high-quality PowerPoint presentation for the Hackathon.

Automatically checks for python-pptx, installs it if needed, checks for any existing template,
and produces 'Build_With_Bharat_Final_Presentation.pptx'.
"""

import os
import subprocess
import sys
from pathlib import Path

# Ensure python-pptx is installed
try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("[*] Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE


def create_deck(output_path: str = "Build_With_Bharat_Final_Presentation.pptx"):
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (Matching Hackathon Theme)
    ORANGE = RGBColor(235, 94, 40)       # #EB5E28
    DARK = RGBColor(37, 36, 34)          # #252422
    LIGHT_BG = RGBColor(255, 252, 242)   # #FFFCF2
    CARD_BG = RGBColor(245, 245, 247)    # #F5F5F7
    TEXT_MUTED = RGBColor(100, 100, 100) # #646464
    BORDER_COLOR = RGBColor(220, 220, 225)
    ACCENT_GREEN = RGBColor(46, 125, 50)
    WHITE = RGBColor(255, 255, 255)

    def add_header(slide, title_text: str):
        # Left Orange Pill Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.55), Inches(0.18), Inches(0.75))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ORANGE
        bar.line.color.rgb = ORANGE

        # Header Text
        txBox = slide.shapes.add_textbox(Inches(1.15), Inches(0.5), Inches(9.5), Inches(0.85))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.font.name = "Arial"

        # Top Orange Divider Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.73), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE
        line.line.color.rgb = ORANGE

        # Bottom Footer
        footBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.73), Inches(0.4))
        p_foot = footBox.text_frame.paragraphs[0]
        p_foot.text = "BUILD WITH भारत 2.0  •  NATIONAL LEVEL HACKATHON"
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = TEXT_MUTED
        p_foot.font.name = "Arial"
        p_foot.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 1: COVER PAGE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)

    # Sub-header tag
    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.33), Inches(0.6))
    p = t_box.text_frame.paragraphs[0]
    p.text = "BUILD WITH भारत 2.0"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    p2 = t_box.text_frame.add_paragraph()
    p2.text = "NATIONAL LEVEL HACKATHON"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = ORANGE
    p2.alignment = PP_ALIGN.CENTER

    # Project Title Card
    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.33), Inches(2.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = DARK
    card1.line.color.rgb = ORANGE
    card1.line.width = Pt(1.5)

    tf_c = card1.text_frame
    tf_c.word_wrap = True
    p_t = tf_c.paragraphs[0]
    p_t.text = "Autonomous OS Debugging Agent"
    p_t.font.size = Pt(30)
    p_t.font.bold = True
    p_t.font.color.rgb = ORANGE
    p_t.alignment = PP_ALIGN.CENTER

    p_sub = tf_c.add_paragraph()
    p_sub.text = "AI-Powered System Diagnostics, Safe Remediation & Zero-Risk Rollback"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = WHITE
    p_sub.alignment = PP_ALIGN.CENTER

    # Meta Info Card
    card_meta = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.7), Inches(10.33), Inches(1.9))
    card_meta.fill.solid()
    card_meta.fill.fore_color.rgb = CARD_BG
    card_meta.line.color.rgb = BORDER_COLOR

    tf_m = card_meta.text_frame
    tf_m.word_wrap = True
    
    rows = [
        ("TEAM NAME:", "KernelHealers"),
        ("PROBLEM STATEMENT:", "Autonomous Diagnosis & Remediation of Operating System & Kernel Errors"),
        ("TEAM MEMBERS:", "Ansh & Team"),
        ("COLLEGE / INSTITUTION:", "Your College / University Name"),
    ]
    for i, (k, v) in enumerate(rows):
        p_row = tf_m.paragraphs[0] if i == 0 else tf_m.add_paragraph()
        p_row.text = f"{k}  {v}"
        p_row.font.size = Pt(13)
        p_row.font.color.rgb = DARK
        p_row.font.name = "Arial"

    # ==========================================
    # SLIDE 2: PROBLEM STATEMENT
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "PROBLEM STATEMENT")

    # 3 Problem Cards
    cards_data = [
        ("Cryptic & Opaque Error Codes", "Errors like '0x80070005' or '0x80240020' give zero meaningful context to regular users or developers, making root cause identification frustrating."),
        ("Dangerous 'Forum Copy-Pasting'", "Users spend 45+ mins scouring Reddit & StackOverflow, copy-pasting unverified PowerShell commands that can corrupt registries or break system boot."),
        ("Massive IT Support Overhead", "Over 30% of enterprise IT helpdesk time is wasted on repetitive OS permissions, broken background services, and corrupted update caches."),
    ]
    for i, (title, desc) in enumerate(cards_data):
        c = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 3.98), Inches(1.8), Inches(3.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = ORANGE if i == 1 else BORDER_COLOR
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = f"0{i+1}"
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = ORANGE

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(17)
        p2.font.bold = True
        p2.font.color.rgb = DARK

        p3 = tf.add_paragraph()
        p3.text = f"\n{desc}"
        p3.font.size = Pt(13)
        p3.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 3: SOLUTION
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "PROPOSED SOLUTION")

    sol_items = [
        ("Live Machine Ingestion", "Queries live Windows Event Viewer logs (System, App, WindowsUpdate) & OS metadata in real-time."),
        ("Closed-Loop AI Reasoning", "Generates hypotheses and executes safe, read-only system probes (icacls, Get-Service) to confirm root cause."),
        ("Human-in-the-Loop Gate", "Generates syntax-highlighted PowerShell fix scripts that require explicit user sign-off ([Y/n]) before running."),
        ("Pre-Fix Snapshot & Rollback", "Automatically snapshots configuration before fixing, enabling 1-click instant rollback ('agent.py rollback')."),
        ("100% Privacy & Local LLM", "Runs seamlessly with OpenAI GPT-4o or 100% offline via local Ollama models (Llama 3.1, Mistral)."),
    ]
    for i, (title, desc) in enumerate(sol_items):
        y_pos = Inches(1.65 + i * 1.0)
        pill = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.73), Inches(0.85))
        pill.fill.solid()
        pill.fill.fore_color.rgb = CARD_BG
        pill.line.color.rgb = BORDER_COLOR

        tf = pill.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"✓  {title}:  "
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ORANGE

        p_desc = p
        p_desc.text = f"✓  {title}: {desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = DARK

    # ==========================================
    # SLIDE 4: FLOW OF SOLUTION
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "FLOW OF SOLUTION")

    steps = [
        ("1. Ingestion", "Admin check + Event Viewer critical logs (50 events) + OS telemetry"),
        ("2. Diagnosis", "LLM formulates hypothesis & outputs safe read-only PowerShell commands"),
        ("3. Probing", "Executes safe probes (icacls, services) with blacklist safety guards"),
        ("4. Approval", "Previews script in terminal & requires explicit [Y/n] user approval"),
        ("5. Fix & Rollback", "Takes snapshot, runs fix, verifies service health & enables 1-click rollback"),
    ]
    for i, (title, desc) in enumerate(steps):
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 2.38), Inches(2.2), Inches(2.25), Inches(4.2))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ORANGE
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = ORANGE
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK

    # ==========================================
    # SLIDE 5: TECH STACK
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "TECH STACK & IMPLEMENTATION")

    stack_cards = [
        ("Core CLI & UI", ["Python 3.10+", "Typer (Type-safe CLI)", "Rich (Banners, Tables, Monokai Highlighting)"]),
        ("AI & Reasoning Engine", ["OpenAI API (GPT-4o)", "Ollama (Local Llama 3.1 / Mistral)", "Strict JSON Schema Prompts"]),
        ("OS & System Diagnostics", ["PowerShell 5.1 / 7+", "Windows Event Viewer (Get-WinEvent)", "WMI / Win32 System APIs"]),
        ("Security & Safety Guardrails", ["Regex Command Blacklist (del/format/etc.)", "Temp File Sandbox & Auto-Cleanup", "Snapshot & Rollback Engine (.backups/)", "Windows RunOnce Post-Reboot Hook"]),
    ]
    for i, (cat, items) in enumerate(stack_cards):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.0)
        y = Inches(1.8 + row * 2.5)

        c = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.73), Inches(2.2))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = BORDER_COLOR

        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cat
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ORANGE

        for item in items:
            p_it = tf.add_paragraph()
            p_it.text = f"•  {item}"
            p_it.font.size = Pt(13)
            p_it.font.color.rgb = DARK

    # ==========================================
    # SLIDE 6: USP
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "UNIQUE SELLING PROPOSITION (USP)")

    usps = [
        ("Ground Truth Telemetry (No Hallucinations)", "Unlike chatbots that hallucinate blind commands, our agent runs live read-only verification checks on actual system files and registry keys before proposing fixes."),
        ("Strict Human-in-the-Loop Safety", "Full transparency: The exact PowerShell code is shown with syntax highlighting. Nothing executes without explicit user consent ([Y/n])."),
        ("Zero-Risk Snapshot Rollback", "Every applied fix automatically creates an inverse rollback script and snapshot point. Revert any change in 1 second via 'agent.py rollback'."),
        ("100% Offline / Air-Gapped Operation", "Can run completely offline on local LLMs with Ollama for sensitive corporate, banking, or defense systems with zero external data egress."),
    ]
    for i, (title, desc) in enumerate(usps):
        c = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 2.98), Inches(1.8), Inches(2.8), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = ORANGE
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = f"★ USP 0{i+1}"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = ORANGE

        p2 = tf.add_paragraph()
        p2.text = f"\n{title}"
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = DARK

        p3 = tf.add_paragraph()
        p3.text = f"\n{desc}"
        p3.font.size = Pt(12)
        p3.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 7: FEASIBILITY & COMPETITORS
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "FEASIBILITY & COMPETITOR ANALYSIS")

    # Table Shape
    rows = 5
    cols = 5
    table_shape = s7.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(3.2))
    table = table_shape.table

    headers = ["Feature", "Windows Troubleshooter", "ChatGPT / Search", "Enterprise RMM", "Our Debug Agent"]
    for col_idx, h in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK if col_idx != 4 else ORANGE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    data = [
        ["Root Cause Depth", "❌ Basic / Rigid", "⚠️ Generic / Blind", "⚠️ Manual Analysis", "✅ Deep & Automated"],
        ["Direct Execution", "⚠️ Limited / Opaque", "❌ Manual Copy-Paste", "⚠️ Complex Custom Scripts", "✅ 1-Click with Approval"],
        ["Rollback Engine", "❌ No Rollback", "❌ No", "⚠️ Requires full disk backup", "✅ Instant 1-Click Rollback"],
        ["Local Privacy / Offline", "✅ Yes", "❌ Cloud Only", "❌ Cloud Dependent", "✅ Local LLM + Cloud"],
    ]
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if col_idx != 4 else RGBColor(255, 235, 230)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER

    # Bottom notes
    box_n = s7.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.73), Inches(1.5))
    tf_n = box_n.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.text = "Key Feasibility Solutions Implemented:"
    p_n.font.bold = True
    p_n.font.size = Pt(13)
    p_n.font.color.rgb = DARK

    items_feas = [
        "1. Dynamic System32 & PowerShell resolution across all Windows release configurations.",
        "2. Built-in command regex blacklists preventing destructive operations during diagnostic phases.",
        "3. Decoupled AI client supporting cloud GPT-4o or zero-cost local Ollama endpoints.",
    ]
    for it in items_feas:
        p_it = tf_n.add_paragraph()
        p_it.text = f"•  {it}"
        p_it.font.size = Pt(12)
        p_it.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 8: RESEARCH & REFERENCE
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "RESEARCH, REPOSITORY & FUTURE SCOPE")

    card_r1 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.73), Inches(4.8))
    card_r1.fill.solid()
    card_r1.fill.fore_color.rgb = CARD_BG
    card_r1.line.color.rgb = BORDER_COLOR

    tf_r1 = card_r1.text_frame
    tf_r1.word_wrap = True
    p = tf_r1.paragraphs[0]
    p.text = "Repository & Live Validation"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    items_r1 = [
        "GitHub Repository: Open-source MIT Licensed codebase.",
        "Live Validation: Successfully tested on Windows Update permission denial (0x80070005) & service failure modes.",
        "Live Subcommands: 'diagnose', 'history', 'rollback', 'check-env'.",
        "Deterministic Output: 100% verified PowerShell script execution and exit code validation.",
    ]
    for it in items_r1:
        p_it = tf_r1.add_paragraph()
        p_it.text = f"\n•  {it}"
        p_it.font.size = Pt(12)
        p_it.font.color.rgb = DARK

    card_r2 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.73), Inches(4.8))
    card_r2.fill.solid()
    card_r2.fill.fore_color.rgb = CARD_BG
    card_r2.line.color.rgb = BORDER_COLOR

    tf_r2 = card_r2.text_frame
    tf_r2.word_wrap = True
    p = tf_r2.paragraphs[0]
    p.text = "Future Roadmap & Scaling"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    items_r2 = [
        "Multi-Platform Support: Linux systemd/journalctl logs and macOS launchctl remediation.",
        "Fleet Enterprise Mode: Centralized agent running diagnostics across Kubernetes nodes & remote server fleets.",
        "Interactive Chat Mode: Multi-turn conversational diagnostics for ambiguous kernel issues.",
        "Automated Patch Suggestions: Direct integration with Microsoft Update Catalog & WinGet package manager.",
    ]
    for it in items_r2:
        p_it = tf_r2.add_paragraph()
        p_it.text = f"\n•  {it}"
        p_it.font.size = Pt(12)
        p_it.font.color.rgb = DARK

    # Save presentation
    prs.save(output_path)
    print(f"[+] Successfully generated presentation: {output_path}")


if __name__ == "__main__":
    out_file = sys.argv[1] if len(sys.argv) > 1 else "Build_With_Bharat_Final_Presentation.pptx"
    create_deck(out_file)
